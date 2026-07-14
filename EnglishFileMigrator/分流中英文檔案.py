# ==========================================================
# MODULE:       Script_EnglishFileMigrator
# PURPOSE:      自動掃描指定資料夾，判定語系並安全移轉英文檔案至專屬目錄，同時生成溯源對照表
# EXPORTS:      EnglishFileMigrator
# IMPORTS:      os, shutil, csv, hashlib, datetime, pathlib, pandas, docx, pdfminer, pptx
# FORBIDDEN:    禁止使用 open('w') 直接覆寫正式報表；禁止使用未經驗證的直接移動（shutil.move）
# DEPENDENCIES: ACDS_ContractRegistry (CONFIG 變數結構), ACDS_ADR (ADR-001, ADR-002)
# VERSION:      1.1.0 [Stability: Experimental]
# ==========================================================

import csv
from datetime import datetime
import hashlib
import os
from pathlib import Path
import shutil
import sys

# 第三方套件
from docx import Document
import pandas as pd
from pdfminer.high_level import extract_text
from pptx import Presentation

# ==========================================
# SSOT: 語系自動判定配置 (唯一入口)
# ==========================================

CONFIG = {
    "CHINESE_THRESHOLD": 10,
    "PDF_PAGE_LIMIT": 30,
    "TARGET_EXTS": [
        '.pdf', '.docx', '.txt', '.xlsx', '.xls', '.csv', '.pptx',
        '.jpg', '.jpeg', '.png', '.gif'
    ],
    "REPORT_NAME": "英文檔案移轉對照表.csv",
    "ENGLISH_FOLDER": "英文區"
}


class ContentExtractionError(Exception):
    """自訂例外：內容解析失敗，避免靜默回傳 False"""
    pass


class EnglishFileMigrator:
    def __init__(self, root_dir):
        self.root = Path(root_dir).resolve()
        self.eng_folder = self.root / CONFIG["ENGLISH_FOLDER"]
        self.report_path = self.root / CONFIG["REPORT_NAME"]
        self.temp_report_path = self.root / f"{CONFIG['REPORT_NAME']}.tmp"
        self.stats = {"moved": 0, "stayed": 0, "errors": 0}
        self.module_name = "英文移轉模組"

    def _log(self, level, message):
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        print(f"[{level}] {timestamp} {self.module_name} {message}")

    def _get_file_hash(self, path):
        """計算檔案 SHA-256 雜湊值"""
        hasher = hashlib.sha256()       
        try:
            with open(path, 'rb') as f:
                for chunk in iter(lambda: f.read(65536), b''):
                    hasher.update(chunk)
            return hasher.hexdigest()
        except PermissionError:
            raise
        except Exception as e:
            self._log("ERROR", f"計算雜湊失敗 {path.name}: {e}")
            return None

    def is_english_content(self, path):
        text = path.stem 
        ext = path.suffix.lower()    

        try:
            if ext == '.pdf':
                text += extract_text(path, page_numbers=list(range(CONFIG['PDF_PAGE_LIMIT'])))
            elif ext == '.docx':
                doc = Document(path)
                text += " ".join([p.text for i, p in enumerate(doc.paragraphs) if i < CONFIG['PDF_PAGE_LIMIT']])
            elif ext == '.txt':
                try:
                    text += path.read_text(encoding='utf-8')
                except Exception:
                    text += path.read_text(encoding='big5', errors='ignore')
            elif ext in ['.xlsx', '.xls', '.csv']:
                try:
                    df = pd.read_csv(path, nrows=10) if ext == '.csv' else pd.read_excel(path, nrows=10)
                    text += " ".join(df.columns.astype(str))
                except Exception as e:
                    raise ContentExtractionError(str(e))
            elif ext == '.pptx':
                try:
                    prs = Presentation(path)
                    for slide in prs.slides[:CONFIG['PDF_PAGE_LIMIT']]:
                        text += " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                except Exception as e:
                    raise ContentExtractionError(str(e))
        except PermissionError:
            raise
        except Exception as e:
            # 統一捕捉未被內層攔截的非預期錯誤
            raise ContentExtractionError(str(e))

        chinese_count = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
      
        if ext in ['.jpg', '.jpeg', '.png', '.gif']:
            chinese_count += sum(1 for c in path.stem if '\u4e00' <= c <= '\u9fff')
        return chinese_count < CONFIG['CHINESE_THRESHOLD']
        
    def execute(self):
        if not self.root.is_dir():
            self._log("CRIT", f"目標路徑無效: {self.root}")
            return

        self.eng_folder.mkdir(exist_ok=True)     

        try:
            with open(self.temp_report_path, 'w', newline='', encoding='utf-8-sig') as f:
                writer = csv.writer(f)
                writer.writerow(["檔案名稱", "判定結果", "動作狀態", "原始位置(溯源用)", "目前位置"])
                
                for path in self.root.rglob("*"):
                    if not path.is_file():
                        continue
                    if path.suffix.lower() not in CONFIG["TARGET_EXTS"]:
                        continue
                    if CONFIG["ENGLISH_FOLDER"] in path.parts:
                        continue
                    if path.name == CONFIG["REPORT_NAME"] or path.name == self.temp_report_path.name:
                        continue

                    old_full_path = str(path)
                    
                    try:
                        is_eng = self.is_english_content(path)
                    except PermissionError:
                        self._log("WARN", f"檔案遭鎖定，放棄讀取判定: {path.name}")
                        writer.writerow([path.name, "未判定", "檔案鎖定，保留原位", old_full_path, old_full_path])
                        self.stats["stayed"] += 1
                        continue
                    except ContentExtractionError as e:
                        self._log("WARN", f"檔案內文提取失敗，放棄判定: {path.name} (原因: {e})")
                        writer.writerow([path.name, "未判定", "內容解析失敗，保留原位", old_full_path, old_full_path])
                        self.stats["stayed"] += 1
                        self.stats["errors"] += 1
                        continue
                    
                    if is_eng and path.suffix.lower() not in ['.jpg', '.jpeg', '.png', '.gif']:
                        dest = self.eng_folder / path.name
                        
                        if dest.exists():
                            try:
                                src_hash = self._get_file_hash(path)
                                dest_hash = self._get_file_hash(dest)
                            except PermissionError:
                                self._log("WARN", f"檔案遭鎖定，放棄雜湊比對: {path.name}")
                                writer.writerow([path.name, "英文", "檔案鎖定，保留原位", old_full_path, old_full_path])
                                self.stats["stayed"] += 1
                                continue
                            
                            if src_hash is not None and src_hash == dest_hash:
                                action = "重複檔案（雜湊相同），已安全清理"
                                current_pos = str(dest)
                                try:
                                    path.unlink()
                                    self.stats["stayed"] += 1
                                except PermissionError:
                                    action = "清理重複檔案失敗: 權限不足或檔案佔用"
                                    self._log("WARN", action)
                                    self.stats["errors"] += 1
                                writer.writerow([path.name, "英文", action, old_full_path, current_pos])
                                continue
                            else:
                                counter = 1
                                while dest.exists():
                                    dest = self.eng_folder / f"{path.stem}_{counter}{path.suffix}"
                                    counter += 1

                        try:
                            shutil.copy2(path, dest)
                            
                            h1 = self._get_file_hash(path)
                            h2 = self._get_file_hash(dest)
                            
                            if h1 is not None and h1 == h2:
                                path.unlink()
                                action = "已移轉至英文區"
                                current_pos = str(dest)
                                self.stats["moved"] += 1
                                self._log("INFO", f"成功安全移轉檔案: {path.name}")
                            else:
                                raise IOError("複製後雜湊值不一致或無法驗證，已清除殘留副本")
                                
                        except PermissionError:
                            action = "移轉失敗: 檔案遭系統鎖定 (In-Use)"
                            current_pos = old_full_path
                            self.stats["errors"] += 1
                            self._log("WARN", f"檔案佔用防禦觸發: {path.name}")
                            if dest.exists() and str(dest) != str(path):
                                try: 
                                    dest.unlink(missing_ok=True)
                                except Exception as clean_err:
                                    self._log("WARN", f"清理鎖定殘留檔失敗 {dest.name}: {clean_err}")
                                    
                        except Exception as e:
                            action = f"安全移轉失敗: {e}"
                            current_pos = old_full_path
                            self.stats["errors"] += 1
                            self._log("ERROR", f"移轉異常 {path.name}: {e}")
                            if dest.exists() and str(dest) != str(path):
                                try: 
                                    dest.unlink(missing_ok=True)
                                except Exception as clean_err:
                                    self._log("WARN", f"清理異常殘留檔失敗 {dest.name}: {clean_err}")
                    else:
                        action = "保留原位"
                        current_pos = old_full_path
                        self.stats["stayed"] += 1
                        self._log("INFO", f"檔案保留原位: {path.name}")

                    writer.writerow([path.name, "英文" if is_eng else "非英文", action, old_full_path, current_pos])

            self.temp_report_path.replace(self.report_path)
            self._show_summary()

        except (Exception, KeyboardInterrupt) as e:
            err_msg = "使用者強制中斷" if isinstance(e, KeyboardInterrupt) else f"批次任務崩潰: {e}"
            self._log("CRIT", err_msg)
            
            if self.temp_report_path.exists():
                crash_log_path = self.root / f"{CONFIG['REPORT_NAME']}.crashed"
                try:
                    self.temp_report_path.rename(crash_log_path)
                    self._log("CRIT", f"為保留溯源軌跡，未完成之報表已另存為: {crash_log_path.name}")
                except Exception as rename_err:
                    self._log("CRIT", f"無法保留崩潰報表: {rename_err}")
                    
            if isinstance(e, KeyboardInterrupt):
                raise  

    def _show_summary(self):
        self._log("INFO", "-" * 50)
        self._log("INFO", "移轉任務執行結束")
        self._log("INFO", f"統計: 已移動 {self.stats['moved']} | 保留原位 {self.stats['stayed']} | 異常 {self.stats['errors']}")
        self._log("INFO", "（註：發生異常的檔案皆已安全保留原位，因此「保留原位」包含了「異常」的數量）")
        self._log("INFO", f"溯源報表路徑: {self.report_path}")
        self._log("INFO", "-" * 50)


if __name__ == "__main__":
    os.system('cls' if os.name == 'nt' else 'clear')
    
    migrator_logger = EnglishFileMigrator(".")
    migrator_logger._log("INFO", "英文檔案自動化移轉系統啟動")
    print("\n")
    
    try:
        path_input = input("請輸入目標資料夾路徑: ").strip().strip('"')
        if not path_input:
            sys.exit()
            
        migrator = EnglishFileMigrator(path_input)
        migrator.execute()
    except KeyboardInterrupt:
        print("\n")
        migrator_logger._log("WARN", "使用者強制中斷主程式")
    
    print("\n")
    input("任務結束，按 Enter 鍵關閉視窗...")
